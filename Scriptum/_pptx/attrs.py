#!/usr/bin/env python3
# coding: utf-8
#
# part of:
#   S C R I P T U M

#
from typing import Any, Tuple

from pptx.util import Pt
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.oxml.ns import qn
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR

class ShapeAttrs:
    """extract all attributes from the shapes if possible
    this includes fonts and colors"""
    def __init__(self, shape: Any):
        
        # basic settings/sizes
        self.top = shape.top
        self.left = shape.left
        self.width = shape.width
        self.height = shape.height
        self.x = self.left+0.5*self.width
        self.y = self.top +0.5*self.height
        self.shape_type = shape.shape_type
        self.table = None
        self.font = None
        
        if shape.has_text_frame:

            self.font = FontAttrs(shape.text_frame)
        
        if shape.has_table:
            
            self.table = TableAttrs(shape.table)

    def setTableAttrs(self, newtable: Any) -> None:
        if not self.table:
            pass # no overwrite
        else:
            self.table.setAttrs(newtable)
    
    def setFontAttrs(self, newtext_frame: Any) -> None:
        self.font.setAttrs(newtext_frame)


class TableAttrs:
    def __init__(self, table: Any):

        try:
            sid = table._tbl.find(qn('a:tblPr')).find(qn('a:tableStyleId')).text
            quiet = True
        except:
            sid = None
            quiet = False
        self.tableStyle = sid

        if sid:
            skipColors = True
        else:
            skipColors = False

        # extract colors from table
        # experimental: alternate colors vertical and horizontal
        # see PPT "Table Design"
        self.horz_banding = table.horz_banding 
        self.vert_banding = table.vert_banding
        self.first_col = table.first_col # first col different to content
        self.first_row = table.first_row # header different to content
        self.last_col = table.last_col # last col different to content
        self.last_row = table.last_row # total row different to content

        # if not self.first_row: # no header
        #     _header_row = _odd_row = []
        # else:
        #     _header_row = []
        #     _odd_row = []
        # _total_row = [] if self.last_row else None
        # _even_row = []

        self.tableDesign = { }
            #'header_row': _header_row, # first, even, odd, last
        #                         'total_row':  _total_row, # first, even, odd, last
        #                         'even_row':   _even_row, # first, even, odd, last
        #                         'odd_row':    _odd_row, # first, even, odd, last
        #                     }
        
        # we could do that very complex but unreadable
        # thus we have
        # OPTION 1: just a small table
        if len(table.rows) < 3 and len(table.columns) < 3:
            attrs00 = CellAttrs(table.cell(0,0), skipColors)
            try:
                attrs01 = CellAttrs(table.cell(0,1), skipColors)
            except:
                attrs01 = attrs00
            try:
                attrs10 = CellAttrs(table.cell(1,0), skipColors)
            except:
                attrs10 = attrs00
            try:
                attrs11 = CellAttrs(table.cell(1,1), skipColors)
            except:
                attrs11 = attrs01

            # in this case we do ignore all last rows and total colums or odd/even options
            self.tableDesign['header_row'] =  [attrs00, attrs01, attrs01, attrs01]
            self.tableDesign['odd_row'] =     [attrs10, attrs11, attrs11, attrs11]
            self.tableDesign['even_row'] =    [attrs10, attrs11, attrs11, attrs11]
            self.tableDesign['total_row'] =   [attrs10, attrs11, attrs11, attrs11]
            if not quiet:
                print('NOTE: the size of the table template is too small, not much formating extracted!')

        if len(table.rows) >= 3 and len(table.columns) >=3 :
            
            attrs00 = CellAttrs(table.cell(0,0), skipColors)
            attrs01 = CellAttrs(table.cell(0,1), skipColors)
            attrs02 = CellAttrs(table.cell(0,2), skipColors)
            attrs10 = CellAttrs(table.cell(1,0), skipColors)
            attrs11 = CellAttrs(table.cell(1,1), skipColors)
            attrs12 = CellAttrs(table.cell(1,2), skipColors)
            attrs20 = CellAttrs(table.cell(2,0), skipColors)
            attrs21 = CellAttrs(table.cell(2,1), skipColors)
            attrs22 = CellAttrs(table.cell(2,2), skipColors)

            try:
                attrs03 = CellAttrs(table.cell(0,3), skipColors)
                attrs13 = CellAttrs(table.cell(1,3), skipColors)
                attrs23 = CellAttrs(table.cell(2,3), skipColors)
            except:
                # we do ignore last cols
                attrs03 = attrs13 = attrs23 = None 

            try:
                attrs30 = CellAttrs(table.cell(3,0), skipColors)
                attrs31 = CellAttrs(table.cell(3,1), skipColors)
                attrs32 = CellAttrs(table.cell(3,2), skipColors)
            except:
                # we do ignore total row
                attrs30 = attrs31 = attrs32 = None 

            if attrs30 and attrs03:
                attrs33 = CellAttrs(table.cell(3,3), skipColors)
            else:
                attrs33 = None

            self.tableDesign['header_row'] =  [attrs00, attrs01, attrs02, attrs03]
            self.tableDesign['odd_row'] =     [attrs10, attrs11, attrs12, attrs13]
            self.tableDesign['even_row'] =    [attrs20, attrs21, attrs22, attrs23]
            self.tableDesign['total_row'] =   [attrs30, attrs31, attrs32, attrs33]

        if (len(table.rows) > 4 or len(table.columns) > 4) and not quiet :
            print('NOTE: the size of the table template is bigger than required, not all formating used!')

    def setAttrs(self, newtable: Any) -> None:
        
        if self.tableStyle:
            try:
                newtable._tbl.find(qn('a:tblPr')).find(qn('a:tableStyleId')).text = self.tableStyle
                
            except:
                pass

        last_row = len(newtable.rows)-1
        last_col = len(newtable.columns)-1
        for i in range(len(newtable.rows)):
            for j in range(len(newtable.columns)):
                if i == 0 and self.first_row:
                    if j == 0 and self.first_col:
                        attr = self.tableDesign['header_row'][0]
                    elif j == last_col and self.last_col and self.tableDesign['header_row'][-1]:
                        attr = self.tableDesign['header_row'][-1]
                    elif j % 2 == 0:
                        attr = self.tableDesign['header_row'][2]
                    else:
                        attr = self.tableDesign['header_row'][1]
                elif i == last_row and self.last_row and self.tableDesign['total_row'][0]:
                    if j == 0 and self.first_col:
                        attr = self.tableDesign['total_row'][0]
                    elif j == last_col and self.last_col and self.tableDesign['total_row'][-1]:
                        attr = self.tableDesign['total_row'][-1]
                    elif j % 2 == 0:
                        attr = self.tableDesign['total_row'][2]
                    else:
                        attr = self.tableDesign['total_row'][1]
                else:
                    if i % 2 == 0:
                        if j == 0 and self.first_col:
                            attr = self.tableDesign['even_row'][0]
                        elif j == last_col and self.last_col and self.tableDesign['even_row'][-1]:
                            attr = self.tableDesign['even_row'][-1]
                        elif j % 2 == 0:
                            attr = self.tableDesign['even_row'][2]
                        else:
                            attr = self.tableDesign['even_row'][1]
                    else:
                        if j == 0 and self.first_col:
                            attr = self.tableDesign['odd_row'][0]
                        elif j == last_col and self.last_col and self.tableDesign['odd_row'][-1]:
                            attr = self.tableDesign['odd_row'][-1]
                        elif j % 2 == 0:
                            attr = self.tableDesign['odd_row'][2]
                        else:
                            attr = self.tableDesign['odd_row'][1]

                attr.setAttrs(newtable.cell(i,j))
                    



class CellAttrs:
    def __init__(self, cell: Any, skipColors: bool):
        if skipColors:
            self.color = None
        else:
            self.color = getColorFromSolidFill(cell._tc.tcPr.solidFill)
        self.font = FontAttrs(cell.text_frame)
        #
        # margins
        self.margin_left = cell.margin_left
        self.margin_top = cell.margin_top
        self.margin_right = cell.margin_right
        self.margin_bottom = cell.margin_bottom


    def setAttrs(self, newCell: Any) -> None:
        if self.color:
            setCellBackground(newCell, self.color)
        self.font.setAttrs(newCell.text_frame)

        newCell.margin_left = self.margin_left
        newCell.margin_top = self.margin_top
        newCell.margin_right = self.margin_right
        newCell.margin_bottom = self.margin_bottom

class FontAttrs:
    def __init__(self, text_frame: Any):
        
        # this is for.. what?
        self.margin_left = text_frame.margin_left
        self.margin_top = text_frame.margin_top
        self.margin_right = text_frame.margin_right
        self.margin_bottom = text_frame.margin_bottom
        self.vertical_anchor = text_frame.vertical_anchor

        # some attrs are stored in the runs not in the paragraphs...
        p0 = text_frame.paragraphs[0]
        # some attrs are stored in the runs not in the paragraphs...
        if len(p0.runs) >= 1:
            r0 = p0.runs[0]
            #p0.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
            self.alignment = p0.alignment
            self.font_size = r0.font.size
            self.font_italic = r0.font.italic
            self.font_bold = r0.font.bold
            self.font_underline = r0.font.underline
            self.font_name = r0.font.name
            self.font_color = r0.font.color
            #print('GETF',self.font_name, p0.font.bold, p0.text)
        else:
            self.alignment = None
            self.font_size = None # do not set if not known
            self.font_italic = False
            self.font_bold = False
            self.font_underline = False
            self.font_name = None # inherit
            self.font_color = None

        # #defaults
        # self.margin_left = Pt(5)
        # self.margin_top = Pt(8)
        # self.margin_right = Pt(5)
        # self.margin_bottom = Pt(5)
        # self.vertical_anchor = None
        # #p0.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

    def setAttrs(self, newtext_frame: Any) -> None:
        """set attrs to a new shape"""
        newtext_frame.margin_left = self.margin_left
        newtext_frame.margin_top = self.margin_top
        newtext_frame.margin_right = self.margin_right
        newtext_frame.margin_bottom = self.margin_bottom
        newtext_frame.vertical_anchor = self.vertical_anchor
        p0 = newtext_frame.paragraphs[0]
        #p0.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
        p0.alignment = self.alignment
        if self.font_size is not None:
            p0.font.size = self.font_size
        p0.font.italic = self.font_italic
        p0.font.bold = self.font_bold
        p0.font.underline = self.font_underline
        p0.font.name = self.font_name
        if self.font_color:
            if self.font_color.type == MSO_COLOR_TYPE.RGB:
                p0.font.color.rgb = self.font_color.rgb
            elif self.font_color.type == MSO_COLOR_TYPE.SCHEME:
                p0.font.color.theme_color = self.font_color.theme_color
            else:
                pass
        #print('APPLYFONT',p0.font.bold, p0.font.underline, p0.font.size, p0.font.name)
        

def getColorFromSolidFill(solidFill: Any) -> Tuple[str, Any]:
    """
    tested for pptx only"""
    try:
        if solidFill.srgbClr is not None:
            color = ('rgb',RGBColor.from_string(solidFill.srgbClr.values()[0]))
        elif solidFill.schemeClr is not None:
            color = ('theme',_mapSchemeColor(solidFill.schemeClr.values()[0]))
        else:
            color = ('unknown','')
    except Exception as e:
        color = ('failed', f'Reason: {e}')

    return color
    
def _mapSchemeColor(xml_color: str) -> Any:
    """try to reverse map from MSO_THEME_COLOR
    tested for pptx only"""
    for col in MSO_THEME_COLOR:
        if col.xml_value == xml_color:
            #print(col.value)
            #return col.value
            return col
    else:
        return MSO_THEME_COLOR['NOT_THEME_COLOR']
        
def setCellBackground(cell: Any, color: Tuple[str, Any]) -> None:
    """set a color given by a name for a cell
    tested for pptx only"""
    
    ctype, colorval = color
    if ctype == 'rgb':
        cell.fill.solid()
        cell.fill.fore_color.rgb = colorval
    elif ctype == 'theme':
        cell.fill.solid()
        cell.fill.fore_color.theme_color = colorval
    else:
        cell.fill.background() # replace by no color 
        


def UNUSED_setFontColor(run: Any, color: Tuple[str, Any]) -> None:
    """set a color given by a name for a given run
    tested for pptx only"""
    
    ctype, colorval = color
    #print(ctype,type(colorval))
    if ctype == 'rgb':
        run.font.color.rgb = colorval
    elif ctype == 'theme':
        run.font.color.theme_color = colorval
    else:
        pass # do nothing
        
