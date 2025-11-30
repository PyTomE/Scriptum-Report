#!/usr/bin/env python3
# coding: utf-8
#
# part of:
#   S C R I P T U M 
#
##################
# MODULE _pptx.reportPptx
# PROVIDES 
#   class ManagedPptx - manage a full powerpoint pptx-file and combines all other modules in the package
#                       main entry for powerpoint-documents


#
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from .units import units
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement

import os
if os.name == 'nt':
    import win32com.client

#from .pptElement import PptElement
from .base import genericFill
from .base import extractFontAndDecorators
from .attrs import getColorFromSolidFill
from .templates import makeTemplate
from ..tag.tag import getTag
from .shapes import getShapes
from .slide import Slide
from .. import version

from ..rdf.tasks.report_task import ReportTask

# and what not
#_NOT_ALLOWED = '\\,:;~+*#&%$' # by default OPENING and CLOSING will be added to this list, never tested for unicode characters
#_WRONG_CLOSING = '>'

class ManagedPptx:
    """manage the template given by a presentation
    first step: extract everything
    second step: use these values (tbd)
    """
    def __init__(self, document_name, debug=False):
        """open the presentations and initialize several things"""
        self.document_name = document_name
        self._debug=debug
        
        if debug:
            print('Loading presentation template...')
        self.document = Presentation(document_name)

        # first of all we scan all layouts for syntactical errors
        # and
        # sort them


        # to be filled:
        self.layouts = {} # this will hold all layouts that can be used for new slides
        self.config_layouts = {} # this will hold all layouts named "config:foo" that contain configs for colors, sizes etc.
        self.template_layouts = {} # this will hold all layouts named "template:bar" that contain structures to be used elsewhere

        # first of all scan it
        
        errors = self.extractLayouts()
        self.errors = errors
        self.warnings = []

        if errors:
            print('\n'.join(errors))
            print('every further error may be a result of these errors...')
        
        # most of the following is experimental
        
        self.extractConfig()
        
        self.extractTemplates()
        
        if debug:
            if errors:
                print('... loaded with errors!')
            else:
                print('... successfully loaded!')

        # store the template in action
        # yes, that is dangerous in case we are working with more than one template in the same session
        self.slides = [] # we ignore all existing slides!
        self.collectglobal = []
        self.allelements = {}
        
    #~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~
    def extractLayouts(self) -> list:
        """extract all layouts in a presentation for further usage
        extract "Config:XXX" layouts as well
        
        check shapes and tags during that loop, but do not store them

        limitation: we do extract only the first master, every other master will be ignored
        """
        
        errors = []
        
        for layout in self.document.slide_layouts:
            lname = layout.name.lower().replace(' ','').strip()

            # we assume that there is no layouts with the same normalizd name!
            if lname in self.config_layouts or lname in self.layouts or lname in self.template_layouts:
                errors += [f"found multiple layouts that return the normalized name {lname!r}, laout is {layout.name!r}"]

            #print(layout.name)
            if lname.startswith("documentation"):
                # ignore documentation - those layouts will never be used
                continue  

            elif lname.startswith("config:"):
                self.config_layouts[lname.replace("config:","")] = layout

                # scan all shapes for errors
                allshapes = getShapes(layout, placeholder=False)
                for shape in allshapes:
                    # tags must be opened and closed in the same shape
                    open = []
                    for tag in shape.tags:
                        if tag.tagtype == 'open':
                            open += [tag.puretag]
                        elif tag.tagtype == 'close' and open and tag.puretag == open[-1]:
                            open = open[:-1]
                    if open:
                        errors += [f"xml element(s) are not fully closed in layout {layout.name!r}"]
                        errors += [f"   {o!r} unclosed" for o in open]

            elif lname.startswith("template:"):
                # layouts which contain layouts
                self.template_layouts[lname.replace("template:","")] = layout
                
                # scan all shapes for errors
                allshapes = getShapes(layout, placeholder=False)
                for shape in allshapes:
                    # tags must be opened and closed in the same shape
                    open = []
                    for tag in shape.tags:
                        if tag.tagtype == 'open':
                            open += [tag.puretag]
                        elif tag.tagtype == 'close' and open and tag.puretag == open[-1]:
                            open = open[:-1]
                    if open:
                        errors += [f"xml element(s) are not fully closed in layout {layout.name!r}"]
                        errors += [f"   {o!r} unclosed" for o in open]
            else:
                # normal pptx layout slides
                self.layouts[lname] = layout
                
                # scan all shapes for errors
                allshapes = getShapes(layout, placeholder=True)
                for element in allshapes:
                    #print(element, element.isplaceholder)
                    for tag in element.tags:
                        #print(tag.ns, tag.name, tag.tagtype)
                        if tag.tagtype != 'simple':
                            errors += [f"Expect simple tag in {tag.rawtag} for layout {layout.name!r}" ]

        return errors

    #~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~
    def extractConfig(self, verbose=False):
        """used to extract some defaults and settings defined in the template
        
        """
        self.config = {}

        # extract colors
        # There is a single table on that layout to store all known colors
        # each cell contains a different color
        colors = {}
        for layoutname in self.config_layouts.keys():
            if not layoutname.lower().startswith('colors'):
                continue
            layout = self.config_layouts[layoutname]
            for shape in layout.shapes:
                if hasattr(shape,'has_table') and shape.has_table:
                    for r in shape.table.rows:
                        for cell in r.cells:
                            if cell.text:
                                colorname = cell.text.lower()
                                solidFill = cell._tc.tcPr.solidFill
                                colors[colorname] = getColorFromSolidFill(solidFill)
        self.config['colors'] = colors

        # extract fonts and sizes
        #    <p type=xxx/> - for normal paragraphs, captions, headers
        #    <p type=code/> - for code font, just font spec, size will be derived by paragraph
        #    <p type=.../> - custom font for special words
        #    <size:tiny/>, <size:small/>, ... - just the font sizes, follow latex or a different notation
        
        pars = {}
        sizes = {}
        for layoutname in self.config_layouts.keys():
            if not layoutname.lower().startswith('defaults'):
                continue

            for shape in layout.shapes:
                if not hasattr(shape,'text'):
                    continue
                tags = getTag(shape.text)
                if tags:
                    # extract only first tag
                    tag = tags[0]
                    if tag.ns == 'p' and tag.name == 'p' and 'type' in tag.args:
                        #print(tag.puretag, tag.args)
                        pars[tag.args['type']] = extractFontAndDecorators(shape.text_frame.paragraphs[0])
                    elif tag.ns == 'size':
                        #print(tag.puretag, tag.name)
                        sizes[tag.name] = int(extractFontAndDecorators(shape.text_frame.paragraphs[0])['sz'])#/100
            
        self.config['paragraphs'] = pars
        self.config['sizes'] = sizes
        
    #~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~
    def extractTemplates(self, verbose=False) -> None:
        """extract templates which are
        
        tables with a <table:foo/> tag in the first cell
        groupshapes with a picture and at least a text-shape with a <image:foo/> tag
        shapes with text and a structure or a simpe tag or series of simple tags
        groupshapes with text and ... 
        shapes with the same name:ns tags like table:default and table:default:description
        """
        self.templates = {}
        for layout in self.template_layouts.values():
            #print(layout)
            candidates = {}
            for shape in layout.shapes:
                shapes = getShapes(shape)

                if len(shapes) > 1:
                    # this is a groupshape, use it as it is, there is no combination method implemented yet
                    self.templates.update(makeTemplate(shapes, ttype = 'group'))
                    continue
                
                if not shapes or not shapes[0].tags:
                    continue

                # otherwise those are single shapes with single tags
                tag = shapes[0].tags[0]
                candidates[(tag.name,tag.ns)] = candidates.get((tag.name,tag.ns),[])+[shapes[0]]
            #print(candidates)
            for key,candidate in candidates.items():
                self.templates.update(makeTemplate(candidate, ttype = key[1]))
                
    #~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~
    def add_slide(self, layoutname: str) -> Slide:
        """add a slide given by layoutname"""
        
        lname = layoutname.lower().replace(' ','').strip()
        layout = self.layouts[lname]
        self.currentslide = Slide(self.document, layout)
        self.slides += [self.currentslide]
        
        return self.currentslide

    #~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~
    def remove_slide(self, index: int) -> None:
        """remove a slide given by index"""
        
        xml_slides = self.document.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[index])
 
    #~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~
    def applyTask(self,task:ReportTask):
        """use a task from reportData-classes and apply the content to the document"""
        
        #print(task._inspect())
        value = task.value
        section = task.path[0]
        path = task.path[1:] 
        target = task.target

        #print(section,path,'T',target)
        
        if section == 'global':
            # find globally all elements with that tag
            # collect them all to complete them at the end of everything
            self.collectglobal += [(target,task)]

        elif task.what == 'copy':
            # every time we go back to highest level we create a new slide
            # thus every task after this one works on self.currentslide!
            slide = self.add_slide(section)
            
            # however we need to collect all elements to do the global tasks later, just before saving
            for element in slide.ph:
                #print('E',element, element.tags)
                for tag in element.tags:
                    self.allelements[tag.puretag] = self.allelements.get(tag.puretag,[])+[element]

        elif task.what == 'add':
            # adding content which does not yet exist requires templates
            # to be added where a @marker:content exists as general shape and target
            # we need to find the task.where first
            mpath = [task.where]
            elems = self.findElements(path=mpath)
            # print('add', elems)
            if not elems:
                print(f"CHECK: ups, no way to add, there is no task {mpath[-1]} in path {'.'.join(mpath[:-1])}")
            else:
                if len(elems) > 1:
                    # this should usually not happen - but who knows
                    print("WARNING: found more than one element as base shape for 'add' operation!")
                    print(f"         will take the first only on {mpath!r}")
                elem = elems[0]
                if target in self.templates.keys():
                    tmpl = self.templates[target]
                    tmpl.copyAndFill(self.currentslide.slide, elem.thing, value, task.actions)
                    
                else:
                    print(f"WARNING: ups, there is no template {target} in templates")

        elif value.type == 'numbering':
            print(f'do a numbering {value} {section} {path} {target}')
            for shape in self.currentslide.layout.placeholders:
                print(f"{shape.text} {getTag(shape.text)[0].name}")

            templates = [ self.makeTemplate(shape) 
                          for shape in self.currentslide.layout.placeholders
                          if getTag(shape.text)[0].name == target]
            elements = self.currentslide.multiply_placeholder('tocn',6,templates)
            #for element in elements:
            #    for tag in element.tags:
            #        self.allelements[tag.puretag] = self.allelements.get(tag.puretag,[])+[element]
        elif value.type == 'information only':
            pass # is never used - just for infomation and debugging added 
        else:
            # those are the simple tasks where we simply replace the content/tag
            elems = self.findElements(path=[target])
            #print('###',target,elems)
            if elems:
                genericFill(elems,task)
            
    #~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~
    def findElements(self,inall=None, path=None, start=None, alternative=None):
        """implement different ways to find elements/shapes"""
        #print('fE',inall,path,start,alternative)
        #print(self.currentslide)
        elems = []
        if inall is not None:
            elems += self.allelements.get(inall,[])
        elif path and start is None:
            slide = self.currentslide
            if len(path) > 1:
                print(f"WARNING: confusing path used to address placeholder: {'.'.join(path)}")
                print('         take last section of that path!')
            p = path[-1] # should be one element in all cases
            #print(slide.tag_in_ph)
            elems = slide.tag_in_ph.get(p,[])
        #print('E found', elems)       
        return elems

    #~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~
    def artist(self, rdf, 
               directfill=True,
               globalfill=True,
               cleardust=True,
               setproperties=True):
        """the final marriage between document and rdf"""
        print('painting the shapes:')
        if not directfill:                
            print('   SKIP: fill the canvas...')
        else:
            print('   fill the canvas...')
            for task in rdf.tasks:
                # print(task, task.value, task.value.type)
                self.applyTask(task)
        
        if not globalfill:                
            print('   SKIP: fill the global canvas...')
        else:
            print('   fill the global canvas...')
            if self.collectglobal:
                #print(self.collectglobal)
                # do now the global tasks
                for where, task in self.collectglobal:
                    elems = self.findElements(inall=where)
                    if elems:
                        genericFill(elems, task)
        
        if not cleardust:                
            print('   SKIP: clearing all the dust...')
        else:
            print('   clearing all the dust...')
            # hide all placeholders with a blank
            for elems in self.allelements.values():
                for elem in elems:
                    if elem.thing.has_text_frame and elem.isplaceholder and not elem.thing.text:
                        elem.thing.text = ' '

        if not setproperties:                
            print('   SKIP: set properties...')
        else:
            print('   set properties...')

            documenttitle = rdf.settings.documenttitle
            self.document.core_properties.author = f'Scriptum {version}'
            self.document.core_properties.title = documenttitle

        print('done')
        
    #~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~
    def save(self, filename, finish=False, createpdf=False):
        """do a final cleanup and save the result"""
                            
        # save from with python-pptx
        self.document.save(filename)
        
        if finish:
            if os.name == 'nt':
                in_file = os.path.abspath(filename)
                out_file = os.path.abspath(os.path.splitext(os.path.basename(filename))[0]+'.pdf')
                ppSaveAsPDF = 32

                try:
                    try:
                        powerp = win32com.client.GetActiveObject('PowerPoint.Application')
                        doquit = False
                    except:
                        powerp = win32com.client.Dispatch('PowerPoint.Application')
                        doquit = True
                    #print('Visbility:',powerp.Visible,powerp.Version)
                    try:
                        #When an application is launched by the user, the Visible and UserControl properties 
                        #of the Application object are both set to True. When the UserControl property is set to True, 
                        #it isn't possible to set the Visible property of the object to False.
                        powerp.Visible = False # not allowed/working for PPT!???
                    except:
                        pass
                    ppt = powerp.Presentations.Open(in_file)
                    if createpdf:
                        ppt.SaveAs(out_file, FileFormat=ppSaveAsPDF)
                    ppt.Close(SaveChanges=True)
                    if doquit:
                        powerp.Quit()
                except Exception as e:
                    print(f'failed to update and save as PDF\nReason:\n{e}')

            else:
                print(f'Running on {os.name} will prevent any finishing work...')
